# Generates a PowerPoint deck explaining the Fake News Detection website's features and usage.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "reports", "figures")

# ---- Palette ----
NAVY = RGBColor(0x11, 0x1E, 0x3A)
BLUE = RGBColor(0x2E, 0x5C, 0xF6)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5A, 0x63, 0x76)
GREEN = RGBColor(0x1F, 0xA0, 0x6B)
RED = RGBColor(0xD9, 0x3B, 0x3B)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = color
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=DARK_TEXT, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TEXT, font="Calibri",
                 bullet_color=BLUE, space_after=10, bold_lead=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            lead, rest = item
        else:
            lead, rest = None, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        r0 = p.add_run()
        r0.text = "●  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        r0.font.name = font
        if lead:
            rl = p.add_run()
            rl.text = lead + "  "
            rl.font.size = Pt(size)
            rl.font.bold = True
            rl.font.color.rgb = NAVY
            rl.font.name = font
        r1 = p.add_run()
        r1.text = rest
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = font
    return tb


def add_pagenum(slide, n):
    add_text(slide, SW - Inches(1.0), SH - Inches(0.5), Inches(0.7), Inches(0.35),
              str(n), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle=None):
    fill_bg(slide, WHITE)
    add_rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    add_rect(slide, 0, Inches(1.15), SW, Pt(4), BLUE)
    add_text(slide, Inches(0.55), Inches(0.18), SW - Inches(1.1), Inches(0.7),
              title, size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.28), SW - Inches(1.1), Inches(0.4),
                  subtitle, size=14, color=GRAY, italic=True)


# ============================================================= SLIDE 1 — TITLE
s = add_slide()
fill_bg(s, NAVY)
add_rect(s, 0, Inches(6.7), SW, Inches(0.8), BLUE)
add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
          "Fake News Detection", size=48, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(0.7),
          "A Machine Learning Web App for Classifying News as REAL or FAKE",
          size=22, color=RGBColor(0xC9, 0xD4, 0xF2))
add_text(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.5),
          "Flask  •  TF-IDF  •  Linear SVM  •  93.66% Accuracy", size=16, color=BLUE, bold=True)
add_text(s, Inches(0.9), Inches(6.85), Inches(8), Inches(0.5),
          "Features & User Guide", size=14, color=WHITE, italic=True)

# ============================================================= SLIDE 2 — AGENDA
s = add_slide()
header(s, "Agenda")
agenda = [
    "Project Overview — what this website does and why",
    "Key Features — detector, live predictions, transparency dashboard",
    "How It Works — the ML pipeline behind the scenes",
    "How to Use It — step-by-step walkthrough",
    "Model Results & Performance",
    "API Reference — for developers",
    "Tech Stack",
    "Limitations & Future Improvements",
]
add_bullets(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(5), agenda, size=20, space_after=18)
add_pagenum(s, 2)

# ============================================================= SLIDE 3 — OVERVIEW
s = add_slide()
header(s, "Project Overview", "What problem does this website solve?")
add_text(s, Inches(0.55), Inches(1.75), Inches(6.0), Inches(0.4), "The Problem", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.55), Inches(2.2), Inches(6.0), Inches(3.5), [
    "Misinformation spreads faster than fact-checkers can respond.",
    "Readers need a fast, first-pass signal on article credibility.",
    "This app offers an explainable, local, ML-based triage tool.",
], size=15)

add_text(s, Inches(6.9), Inches(1.75), Inches(5.9), Inches(0.4), "The Solution", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(6.9), Inches(2.2), Inches(5.9), Inches(3.5), [
    "A full-stack ML app: paste a headline/article, get an instant REAL / FAKE label with a confidence score.",
    "Runs entirely on localhost — no external API calls, no cloud dependency, no data leaves your machine.",
    "Trained & evaluated on 6,300+ real-world labelled news articles.",
], size=15)

add_rect(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.85), LIGHT)
add_text(s, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.65),
          "Best model: Linear SVM  —  93.66% accuracy  •  0.982 ROC-AUC on held-out test data",
          size=16, bold=True, color=BLUE, anchor=MSO_ANCHOR.MIDDLE)
add_pagenum(s, 3)

# ============================================================= SLIDE 4 — KEY FEATURES
s = add_slide()
header(s, "Key Features")
cards = [
    ("Instant Detector", "Paste a headline and/or article body on the home page and get a REAL/FAKE prediction in real time via AJAX — no page reload."),
    ("Confidence Score", "Every prediction comes with an animated confidence bar showing how certain the model is, plus a cleaned word count."),
    ("\"Try Sample News\" Button", "One-click demo button auto-fills a sample article so new users can instantly see the tool in action."),
    ("Model Transparency Dashboard", "The /about page shows a full model comparison table, dataset statistics, and 6 evaluation charts — total transparency, no black box."),
    ("REST API", "A JSON POST endpoint (/api/predict) exposes the same model for programmatic / developer use."),
    ("Runs Fully Offline", "No cloud calls, no external services — the entire app and model run on your local machine."),
]
x0, y0 = Inches(0.55), Inches(1.65)
cw, ch, gap = Inches(3.95), Inches(1.7), Inches(0.2)
for i, (title, desc) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = x0 + col * (cw + gap)
    y = y0 + row * (ch + gap)
    card = add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.09), ch, BLUE)
    add_text(s, x + Inches(0.25), y + Inches(0.12), cw - Inches(0.4), Inches(0.4), title, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.4), ch - Inches(0.65), desc, size=11.5, color=GRAY)
add_pagenum(s, 4)

# ============================================================= SLIDE 5 — SITE MAP / PAGES
s = add_slide()
header(s, "Website Map", "Two pages, one API")
rows = [
    ("/", "GET", "Detector UI — submit a headline/article for classification"),
    ("/about", "GET", "Model comparison table, dataset stats, evaluation charts"),
    ("/api/predict", "POST", "JSON prediction API (used by the frontend via fetch)"),
    ("/reports/figures/<file>", "GET", "Serves generated evaluation chart images"),
]
tx, ty = Inches(0.7), Inches(1.9)
tw = Inches(11.9)
col_w = [Inches(3.6), Inches(1.5), Inches(6.8)]
rh = Inches(0.85)
headers_row = ["Route", "Method", "Description"]
cx = tx
for j, htxt in enumerate(headers_row):
    add_rect(s, cx, ty, col_w[j], Inches(0.55), NAVY)
    add_text(s, cx + Inches(0.15), ty, col_w[j] - Inches(0.3), Inches(0.55), htxt, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    cx += col_w[j]
for i, (route, method, desc) in enumerate(rows):
    ry = ty + Inches(0.55) + i * rh
    bgc = LIGHT if i % 2 == 0 else WHITE
    cx = tx
    vals = [route, method, desc]
    for j, val in enumerate(vals):
        add_rect(s, cx, ry, col_w[j], rh, bgc)
        col = BLUE if j == 0 else (GREEN if j == 1 else DARK_TEXT)
        bold = j != 2
        add_text(s, cx + Inches(0.15), ry, col_w[j] - Inches(0.3), rh, val, size=13, bold=bold, color=col, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
add_pagenum(s, 5)

# ============================================================= SLIDE 6 — HOW IT WORKS (PIPELINE)
s = add_slide()
header(s, "How It Works", "The ML pipeline behind every prediction")
steps = [
    ("1. Text Preprocessing", "Lowercasing, URL/HTML stripping, punctuation & digit removal. Title is weighted 2x and concatenated with the body."),
    ("2. Feature Extraction", "TF-IDF vectorization: unigrams + bigrams, up to 50,000 features, English stop-words removed."),
    ("3. Model Selection", "4 classifiers trained & compared (Logistic Regression, Passive Aggressive, Linear SVM, Naive Bayes) — best F1 score auto-selected."),
    ("4. Live Inference", "Saved TF-IDF vectorizer + trained model (.joblib) power instant predictions in the web app."),
]
x0 = Inches(0.55)
y0 = Inches(2.1)
bw = Inches(2.95)
bh = Inches(2.4)
gap = Inches(0.15)
for i, (title, desc) in enumerate(steps):
    x = x0 + i * (bw + gap)
    add_rect(s, x, y0, bw, bh, NAVY if i == 2 else LIGHT)
    tcolor = WHITE if i == 2 else NAVY
    dcolor = RGBColor(0xC9, 0xD4, 0xF2) if i == 2 else GRAY
    add_text(s, x + Inches(0.2), y0 + Inches(0.2), bw - Inches(0.4), Inches(0.6), title, size=15, bold=True, color=tcolor)
    add_text(s, x + Inches(0.2), y0 + Inches(0.85), bw - Inches(0.4), bh - Inches(1.0), desc, size=12, color=dcolor)
    if i < len(steps) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw - Inches(0.02), y0 + bh/2 - Inches(0.12), Inches(0.3), Inches(0.24))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = BLUE; arrow.line.fill.background(); arrow.shadow.inherit = False
add_text(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.4), "Model chosen: Linear SVM (highest F1 on held-out test set)", size=14, italic=True, color=BLUE, bold=True)
add_pagenum(s, 6)

# ============================================================= SLIDE 7 — HOW TO USE: STEP 1-2 (Setup)
s = add_slide()
header(s, "How to Use It — Getting Started", "Setup & first run")
add_text(s, Inches(0.55), Inches(1.65), Inches(6.0), Inches(0.4), "Step 1 — Setup", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.55), Inches(2.1), Inches(6.0), Inches(3.2), [
    "Requires Python 3.10+ and pip.",
    "Create & activate a virtual environment (python -m venv venv).",
    "Install dependencies: pip install -r requirements.txt.",
], size=14.5)

add_text(s, Inches(6.9), Inches(1.65), Inches(5.9), Inches(0.4), "Step 2 — Train & Launch", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(6.9), Inches(2.1), Inches(5.9), Inches(3.2), [
    "Train the model once: python -m src.train (skip if models already exist).",
    "Start the app: python run.py.",
    "Open http://127.0.0.1:5000 in your browser.",
], size=14.5)
add_rect(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.9), LIGHT)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7),
          "Tip: the model, vectorizer, metrics, and evaluation charts are all regenerated automatically each time you run training.",
          size=13, italic=True, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_pagenum(s, 7)

# ============================================================= SLIDE 8 — HOW TO USE: DETECTOR PAGE
s = add_slide()
header(s, "How to Use It — the Detector Page", "Home page ( / )")
add_bullets(s, Inches(0.55), Inches(1.75), Inches(7.3), Inches(4.5), [
    ("1. Enter content:", "type/paste a headline and/or the full article text (at least one field is required)."),
    ("2. Or click \"Try Sample News\":", "auto-fills a demo article so you can test instantly."),
    ("3. Click \"Analyze Article\":", "the form submits via fetch() to /api/predict — no page reload."),
    ("4. Read the result:", "a REAL/FAKE badge appears with an animated confidence bar and the cleaned word count."),
    ("5. Interpret responsibly:", "this is a statistical prediction based on writing style — always verify important claims against trusted sources."),
], size=15, space_after=16)

mock = add_rect(s, Inches(8.2), Inches(1.75), Inches(4.55), Inches(4.6), NAVY)
add_text(s, Inches(8.45), Inches(1.95), Inches(4.05), Inches(0.5), "Detector UI (mock)", size=13, bold=True, color=WHITE)
add_rect(s, Inches(8.45), Inches(2.55), Inches(4.05), Inches(0.5), WHITE)
add_text(s, Inches(8.55), Inches(2.55), Inches(3.8), Inches(0.5), "Headline field", size=11, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(8.45), Inches(3.2), Inches(4.05), Inches(1.15), WHITE)
add_text(s, Inches(8.55), Inches(3.2), Inches(3.8), Inches(1.15), "Article text area", size=11, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
btn = add_rect(s, Inches(8.45), Inches(4.55), Inches(1.9), Inches(0.45), BLUE)
add_text(s, Inches(8.45), Inches(4.55), Inches(1.9), Inches(0.45), "Analyze Article", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
btn2 = add_rect(s, Inches(10.45), Inches(4.55), Inches(2.05), Inches(0.45), RGBColor(0x8A,0x94,0xA8))
add_text(s, Inches(10.45), Inches(4.55), Inches(2.05), Inches(0.45), "Try Sample News", size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(8.45), Inches(5.2), Inches(4.05), Inches(0.9), GREEN)
add_text(s, Inches(8.45), Inches(5.2), Inches(4.05), Inches(0.9), "REAL  •  87.4% confidence", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_pagenum(s, 8)

# ============================================================= SLIDE 9 — HOW TO USE: ABOUT PAGE
s = add_slide()
header(s, "How to Use It — the Stats Page", "Model & Dataset Statistics ( /about )")
add_bullets(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(3.6), [
    ("Stat cards:", "total articles, train/test split, best accuracy, and best F1 score at a glance."),
    ("Model comparison table:", "accuracy, precision, recall, F1, and training time for all 4 candidate models, with the winner highlighted."),
    ("Visual evaluation:", "6 charts — model comparison, confusion matrix, ROC curve, class distribution, article length distribution, and top predictive words."),
    ("Pipeline configuration:", "shows the exact TF-IDF settings and model-selection criteria used, for full transparency."),
], size=15, space_after=14)
add_rect(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.9), LIGHT)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7),
          "Use this page to understand WHY the model made a prediction the way it did — it's the transparency layer of the app.",
          size=13, italic=True, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_pagenum(s, 9)

# ============================================================= SLIDE 10 — RESULTS TABLE
s = add_slide()
header(s, "Model Results", "Performance on the held-out test set")
headers_row = ["Model", "Accuracy", "Precision", "Recall", "F1"]
data_rows = [
    ("Logistic Regression", "0.9200", "0.9387", "0.8986", "0.9182", False),
    ("Passive Aggressive", "0.9350", "0.9420", "0.9271", "0.9345", False),
    ("Linear SVM (best)", "0.9366", "0.9480", "0.9239", "0.9358", True),
    ("Multinomial Naive Bayes", "0.8914", "0.8509", "0.9493", "0.8974", False),
]
tx, ty = Inches(0.7), Inches(1.75)
col_w = [Inches(3.6), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)]
rh = Inches(0.6)
cx = tx
for j, htxt in enumerate(headers_row):
    add_rect(s, cx, ty, col_w[j], Inches(0.55), NAVY)
    add_text(s, cx, ty, col_w[j], Inches(0.55), htxt, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += col_w[j]
for i, row in enumerate(data_rows):
    name, acc, prec, rec, f1, best = row
    ry = ty + Inches(0.55) + i * rh
    bgc = RGBColor(0xE4, 0xEE, 0xFF) if best else (LIGHT if i % 2 == 0 else WHITE)
    cx = tx
    vals = [name, acc, prec, rec, f1]
    for j, val in enumerate(vals):
        add_rect(s, cx, ry, col_w[j], rh, bgc)
        add_text(s, cx, ry, col_w[j], rh, val, size=13, bold=(best or j == 0), color=(BLUE if best else DARK_TEXT), align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER), anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[j]
add_text(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.4), "ROC-AUC (best model): 0.9818   |   Dataset: 6,335 articles, near-perfectly balanced (3,164 FAKE / 3,171 REAL)", size=14, bold=True, color=GREEN)

img_path = os.path.join(FIG, "model_comparison.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.9), Inches(5.35), height=Inches(1.9))
add_pagenum(s, 10)

# ============================================================= SLIDE 11 — CHARTS
s = add_slide()
header(s, "Evaluation Charts", "Rendered live on the /about page")
imgs = [
    ("confusion_matrix.png", "Confusion Matrix"),
    ("roc_curve.png", "ROC Curve"),
    ("top_predictive_words.png", "Top Predictive Words"),
]
x0 = Inches(0.5)
y0 = Inches(1.7)
iw = Inches(4.0)
gap = Inches(0.25)
for i, (fname, cap) in enumerate(imgs):
    x = x0 + i * (iw + gap)
    p = os.path.join(FIG, fname)
    if os.path.exists(p):
        s.shapes.add_picture(p, x, y0, width=iw)
    add_text(s, x, Inches(5.55), iw, Inches(0.4), cap, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_pagenum(s, 11)

# ============================================================= SLIDE 12 — API REFERENCE
s = add_slide()
header(s, "API Reference", "For developers — programmatic access")
add_text(s, Inches(0.55), Inches(1.7), Inches(5.6), Inches(0.4), "POST /api/predict", size=18, bold=True, color=BLUE)

add_text(s, Inches(0.55), Inches(2.25), Inches(5.8), Inches(0.35), "Request", size=14, bold=True, color=NAVY)
code_box = add_rect(s, Inches(0.55), Inches(2.65), Inches(5.8), Inches(1.6), NAVY)
add_text(s, Inches(0.75), Inches(2.8), Inches(5.4), Inches(1.3),
          '{\n  "title": "Senate passes new bill",\n  "text": "Full article text here..."\n}',
          size=13, color=RGBColor(0x9C, 0xE8, 0x7B), font="Consolas")

add_text(s, Inches(6.6), Inches(2.25), Inches(6.2), Inches(0.35), "Response", size=14, bold=True, color=NAVY)
code_box2 = add_rect(s, Inches(6.6), Inches(2.65), Inches(6.2), Inches(1.6), NAVY)
add_text(s, Inches(6.8), Inches(2.8), Inches(5.8), Inches(1.3),
          '{\n  "label": "REAL",\n  "confidence": 87.42,\n  "cleaned_word_count": 112\n}',
          size=13, color=RGBColor(0x9C, 0xE8, 0x7B), font="Consolas")

add_bullets(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(2), [
    "At least one of title or text must be non-empty.",
    ("400", "returned for bad/invalid input."),
    ("503", "returned if the model has not been trained yet."),
    "Used internally by the frontend detector page via JavaScript fetch().",
], size=14.5)
add_pagenum(s, 12)

# ============================================================= SLIDE 13 — TECH STACK
s = add_slide()
header(s, "Tech Stack")
stack = [
    ("ML / Data", "Python, scikit-learn, pandas, NumPy, joblib"),
    ("Visualization", "Matplotlib (6 diagnostic charts)"),
    ("Backend", "Flask (application-factory pattern)"),
    ("Frontend", "HTML5, CSS3, vanilla JavaScript (no build step)"),
    ("Testing", "pytest (preprocessing + inference unit tests)"),
]
y = Inches(1.8)
for name, desc in stack:
    add_rect(s, Inches(0.55), y, Inches(3.0), Inches(0.75), NAVY)
    add_text(s, Inches(0.55), y, Inches(3.0), Inches(0.75), name, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(3.65), y, Inches(9.1), Inches(0.75), LIGHT)
    add_text(s, Inches(3.85), y, Inches(8.8), Inches(0.75), desc, size=14, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.9)
add_pagenum(s, 13)

# ============================================================= SLIDE 14 — LIMITATIONS
s = add_slide()
header(s, "Limitations & Disclaimer")
add_bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(3.5), [
    "Learns statistical writing-style patterns — it does not fact-check claims against real-world knowledge or live sources.",
    "Performance reflects the training distribution (English-language news circa the dataset's collection period); may not generalize to satire, non-English text, or novel misinformation styles.",
    "Intended for educational and research purposes only — do not use as a sole source of truth for high-stakes decisions.",
], size=17, space_after=20)
add_rect(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(0.9), RGBColor(0xFD, 0xEC, 0xEC))
add_rect(s, Inches(0.55), Inches(5.3), Inches(0.09), Inches(0.9), RED)
add_text(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(0.7), "Always verify important claims against trusted primary sources.", size=15, bold=True, color=RED, anchor=MSO_ANCHOR.MIDDLE)
add_pagenum(s, 14)

# ============================================================= SLIDE 15 — FUTURE IMPROVEMENTS
s = add_slide()
header(s, "Future Improvements")
add_bullets(s, Inches(0.7), Inches(1.9), Inches(11.8), Inches(4.5), [
    "Add transformer-based models (e.g., DistilBERT) for comparison against classical baselines.",
    "Support batch/CSV upload for bulk article classification.",
    "Add explainability (e.g., LIME/SHAP) directly in the UI to highlight influential words per prediction.",
    "Containerize with Docker for one-command deployment.",
], size=19, space_after=24)
add_pagenum(s, 15)

# ============================================================= SLIDE 16 — CLOSING
s = add_slide()
fill_bg(s, NAVY)
add_rect(s, 0, Inches(6.7), SW, Inches(0.8), BLUE)
add_text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.0), "Thank You", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.6), "http://127.0.0.1:5000", size=22, color=BLUE, bold=True)
add_text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.5), "Questions & Discussion", size=16, color=RGBColor(0xC9, 0xD4, 0xF2), italic=True)

out_path = os.path.join(BASE, "Fake_News_Detection_Website_Overview.pptx")
prs.save(out_path)
print("Saved:", out_path)
